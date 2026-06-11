"""Extract PPTX content to Markdown with image OCR.

Usage:
    python extract_pptx.py "path\to\file.pptx" [output_dir] [options]

Options:
    --no-ocr            Skip OCR for images
    --no-semantic       Skip semantic optimization
    -h, --help          Show this help message
"""

import os
import sys
import re
import zipfile
import shutil
import subprocess
import argparse
from pathlib import Path
from typing import Optional, Dict, List, Tuple


def parse_args():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(
        description="Extract PPTX content to Markdown with image OCR"
    )

    parser.add_argument("pptx_file", help="Path to the PPTX file")
    parser.add_argument("output_dir", nargs="?", default=None,
                        help="Output directory (default: same as input)")
    parser.add_argument("--no-ocr", action="store_true",
                        help="Skip OCR for images")
    parser.add_argument("--no-semantic", action="store_true",
                        help="Skip semantic optimization")

    return parser.parse_args()


# ==================== OCR阶段清理规则 ====================

def clean_ocr_text(text: str) -> str:
    """
    OCR阶段的基础清理：噪声去除、格式修正
    在获取文本时立即处理，避免污染原始数据
    """
    if not text:
        return text

    rules = [
        # 1. 清理多余竖线（OCR常见噪声）
        (r'\|\s*\|', ' '),
        # 2. 清理重复空格
        (r' {2,}', ' '),
        # 3. 清理孤立单字符噪声
        (r'(?<![a-zA-Z])[dDfFgGhHiIjJkKlLmMnNpPqQrRtTvVwWxXyYzZ](?![a-zA-Z])', ''),
        # 4. 清理连续标点
        (r'([^\w\s\n\-=→]){2,}', ' '),
    ]

    cleaned = text
    for pattern, replacement, *flags in rules:
        regex_flags = flags[0] if flags else 0
        cleaned = re.sub(pattern, replacement, cleaned, flags=regex_flags)

    return cleaned.strip()


# ==================== 语义优化阶段 ====================

class SemanticOptimizer:
    """
    语义理解优化器（预留AI接口）
    
    用于在OCR基础清理之上进行更高级的语义理解和优化：
    - 提取关键信息
    - 智能纠错（基于上下文）
    - 内容分类和结构化
    """

    def __init__(self):
        pass

    def optimize_image_ocr(self, ocr_text: str) -> str:
        """优化单张图片的OCR结果"""
        if not ocr_text:
            return ocr_text

        # 过滤无意义的噪声词
        noise_patterns = [
            r'^\s*(TE|wh|we|CAVA|RR)\s*$',
            r'^\s*[DDEEAA\s.]+$',
            r'^\s*[|\\/\-=_~`^\*\+\>\<\?\[\]\(\)]+\s*$',
        ]
        for pattern in noise_patterns:
            if re.match(pattern, ocr_text, re.IGNORECASE):
                return ''

        return ocr_text.strip()


# ==================== 核心功能 ====================

def install_dependencies():
    """Ensure required packages are installed."""
    required_packages = ["markitdown[pptx]", "Pillow", "easyocr", "python-pptx"]
    print("Checking dependencies...")

    for package in required_packages:
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", package, "-q"],
            capture_output=True,
            text=True
        )
        if result.returncode != 0:
            print(f"Warning: Failed to install {package}", file=sys.stderr)


def extract_markitdown(pptx_path: str, output_md: str) -> bool:
    """Extract PPTX text content using markitdown."""
    result = subprocess.run(
        [sys.executable, "-m", "markitdown", pptx_path],
        capture_output=True,
        text=True
    )

    if result.returncode == 0:
        with open(output_md, 'w', encoding='utf-8') as f:
            f.write(result.stdout)
        return True
    else:
        print(f"markitdown error: {result.stderr}", file=sys.stderr)
        return False


def get_slide_images(pptx_path: str) -> Dict[int, List[Tuple[str, int]]]:
    """
    获取每个幻灯片包含的图片信息
    返回: {幻灯片索引: [(图片文件名, 图片大小), ...]}
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        slide_images = {}
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            images = []
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    image = shape.image
                    # 获取图片文件名和大小（用于匹配）
                    image_filename = image.filename
                    image_size = len(image.blob)  # 获取图片数据大小
                    if image_filename:
                        images.append((os.path.basename(image_filename), image_size))
            if images:
                slide_images[slide_idx] = images
        
        return slide_images
    except ImportError:
        print("Warning: python-pptx not found, cannot map images to slides", file=sys.stderr)
        return {}
    except Exception as e:
        print(f"Error reading slide images: {e}", file=sys.stderr)
        return {}


def extract_images_to_output(pptx_path: str, output_dir: str) -> Tuple[Dict[str, str], Dict[int, List[str]]]:
    """
    提取PPTX中的图片到输出目录（改进版：按幻灯片索引命名）
    返回: (原始文件名映射, 幻灯片→图片列表映射)
    """
    images_dir = os.path.join(output_dir, "images")
    
    # 清理旧的images目录
    if os.path.exists(images_dir):
        shutil.rmtree(images_dir, ignore_errors=True)
    
    # 确保父目录存在
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(images_dir, exist_ok=True)

    # 建立原始文件名到输出路径的映射
    image_map = {}
    
    # 建立幻灯片索引到图片列表的映射
    slide_to_images = {}

    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_images = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image'):
                    image = shape.image
                    original_name = image.filename
                    if not original_name:
                        continue
                    
                    # 获取图片数据
                    image_data = image.blob
                    
                    # 按幻灯片索引和图片顺序命名
                    ext = os.path.splitext(original_name)[1].lower()
                    if ext == '.emf':
                        continue
                    
                    # 使用幻灯片索引命名：slide_{slide_idx}_{shape_idx}{ext}
                    output_name = f"slide_{slide_idx}_{shape_idx}{ext}"
                    output_path = os.path.join(images_dir, output_name)
                    
                    with open(output_path, 'wb') as f:
                        f.write(image_data)
                    
                    image_map[original_name] = output_path
                    slide_images.append(output_name)
            
            if slide_images:
                slide_to_images[slide_idx] = slide_images
        
    except ImportError:
        print("Warning: python-pptx not found, falling back to zip extraction", file=sys.stderr)
        # 降级方案：直接从zip提取
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for name in zf.namelist():
                if name.startswith("ppt/media/"):
                    original_name = os.path.basename(name)
                    if not original_name or original_name.strip() == '':
                        continue
                    ext = os.path.splitext(original_name)[1].lower()
                    if ext == '.emf':
                        continue
                    output_path = os.path.join(images_dir, original_name)
                    with open(output_path, 'wb') as f:
                        f.write(zf.read(name))
                    image_map[original_name] = output_path
    
    return image_map, slide_to_images


def fix_image_links(md_path: str, image_map: Dict[str, str], slide_to_images: Dict[int, List[str]] = None):
    """
    修复Markdown中的图片链接，将占位符替换为实际图片路径（改进版）
    :param slide_to_images: 幻灯片索引到图片列表的映射
    """
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 从文件名中提取数字
    def extract_numbers(s: str) -> str:
        return ''.join([c for c in s if c.isdigit()])

    # 如果有幻灯片到图片的映射，按幻灯片索引精确匹配
    if slide_to_images:
        # 按幻灯片分割（markitdown格式：<!-- Slide number: N -->）
        lines = content.split('\n')
        result_lines = []
        current_slide = None
        slide_image_indices = {}  # 记录每个幻灯片当前使用的图片索引
        
        for line in lines:
            # 检测幻灯片标记（<!-- Slide number: N -->）
            slide_match = re.search(r'<!-- Slide number: (\d+) -->', line)
            if slide_match:
                current_slide = int(slide_match.group(1))
                slide_image_indices[current_slide] = 0
                result_lines.append(line)
            elif current_slide and current_slide in slide_to_images:
                # 当前幻灯片有图片，检查是否是图片链接
                img_match = re.search(r'!\[([^\]]*)\]\(([^)]+)\)', line)
                if img_match:
                    alt_text = img_match.group(1)
                    original_path = img_match.group(2)
                    
                    # 获取当前幻灯片的图片列表
                    slide_images = slide_to_images[current_slide]
                    img_idx = slide_image_indices.get(current_slide, 0)
                    
                    if img_idx < len(slide_images):
                        image_name = slide_images[img_idx]
                        slide_image_indices[current_slide] = img_idx + 1
                        # 替换图片链接
                        new_link = f'![{alt_text}](images/{image_name})'
                        line = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', new_link, line)
                
                result_lines.append(line)
            else:
                # 没有幻灯片图片映射，使用旧的文件名匹配方式
                img_match = re.search(r'!\[([^\]]*)\]\(([^)]+)\)', line)
                if img_match and image_map:
                    alt_text = img_match.group(1)
                    original_path = img_match.group(2)
                    original_name = os.path.basename(original_path)
                    
                    found = False
                    matched_path = None
                    original_numbers = extract_numbers(original_name)
                    
                    # 先尝试精确匹配
                    for src_name, dest_path in image_map.items():
                        if src_name.lower() == original_name.lower():
                            matched_path = dest_path
                            found = True
                            break
                    
                    # 如果精确匹配失败，尝试数字匹配
                    if not found and original_numbers:
                        for src_name, dest_path in image_map.items():
                            src_numbers = extract_numbers(src_name)
                            if src_numbers == original_numbers:
                                matched_path = dest_path
                                found = True
                                break
                    
                    if found and matched_path:
                        rel_path = os.path.relpath(matched_path, os.path.dirname(md_path))
                        rel_path = rel_path.replace('\\', '/')
                        new_link = f'![{alt_text}]({rel_path})'
                        line = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', new_link, line)
                
                result_lines.append(line)
        
        final_content = '\n'.join(result_lines)
    else:
        # 没有幻灯片映射，使用旧的文件名匹配方式
        def replace_image_link(match):
            alt_text = match.group(1)
            original_path = match.group(2)
            original_name = os.path.basename(original_path)
            
            found = False
            matched_path = None
            original_numbers = extract_numbers(original_name)
            
            for src_name, dest_path in image_map.items():
                if src_name.lower() == original_name.lower():
                    matched_path = dest_path
                    found = True
                    break
            
            if not found and original_numbers:
                for src_name, dest_path in image_map.items():
                    src_numbers = extract_numbers(src_name)
                    if src_numbers == original_numbers:
                        matched_path = dest_path
                        found = True
                        break
            
            if found and matched_path:
                rel_path = os.path.relpath(matched_path, os.path.dirname(md_path))
                rel_path = rel_path.replace('\\', '/')
                return f'![{alt_text}]({rel_path})'
            else:
                return f'![{alt_text}]({original_path})'
        
        final_content = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', replace_image_link, content)
    
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(final_content)


def ocr_image(img_path: str) -> Optional[str]:
    """Perform OCR on a single image using easyocr."""
    try:
        import easyocr

        # 初始化easyocr阅读器（中文+英文）
        reader = easyocr.Reader(['ch_sim', 'en'], verbose=False)

        ext = os.path.splitext(img_path)[1].lower()
        if ext == '.emf':
            return None

        # 使用easyocr识别图片
        result = reader.readtext(img_path)
        
        # 提取所有识别到的文本并拼接
        text = ' '.join([item[1] for item in result])

        if text.strip():
            # OCR阶段立即进行基础清理
            cleaned_text = clean_ocr_text(text.strip())
            if cleaned_text:
                return cleaned_text
            return text.strip()
        return None

    except ImportError:
        print("Warning: easyocr not found", file=sys.stderr)
        return None
    except Exception as e:
        print(f"OCR error for {img_path}: {e}", file=sys.stderr)
        return None


def insert_ocr_results(md_path: str, ocr_results: Dict[str, str], 
                       slide_images: Dict[int, List[Tuple[str, int]]], 
                       image_info_map: Dict[str, Tuple[str, int]],
                       optimizer: SemanticOptimizer):
    """
    将OCR结果插入到对应的幻灯片位置（改进版）
    """
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 将内容按幻灯片分割
    slides = re.split(r'(<!-- Slide number: (\d+) -->)', content)
    
    # 重新构建内容，在每个幻灯片后插入对应的OCR结果
    result_parts = []
    i = 0
    while i < len(slides):
        if i + 2 < len(slides) and slides[i].startswith('<!-- Slide number:'):
            # 这是幻灯片标记
            slide_num = int(slides[i+1])
            result_parts.append(slides[i])  # 标记
            result_parts.append(slides[i+1])  # 数字
            
            # 添加幻灯片内容
            i += 2
            if i < len(slides):
                result_parts.append(slides[i])
                i += 1
            
            # 检查是否有该幻灯片的图片OCR结果
            if slide_num in slide_images:
                slide_ocr = []
                for image_name, expected_size in slide_images[slide_num]:
                    # 查找对应的OCR结果
                    found = False
                    matched_ocr_name = None
                    
                    # 尝试按文件名匹配（支持新的slide_{idx}_{shape}.ext格式）
                    for ocr_img_name, ocr_text in ocr_results.items():
                        if image_name.lower() == ocr_img_name.lower():
                            matched_ocr_name = ocr_img_name
                            found = True
                            break
                    
                    # 如果文件名不匹配，尝试按幻灯片索引匹配（新格式）
                    if not found:
                        # 检查是否是slide_N格式的图片
                        slide_pattern = re.compile(r'slide_(\d+)_')
                        ocr_match = slide_pattern.match(image_name)
                        if ocr_match and int(ocr_match.group(1)) == slide_num:
                            # 在ocr_results中查找相同幻灯片的图片
                            for ocr_img_name, ocr_text in ocr_results.items():
                                if ocr_img_name.startswith(f'slide_{slide_num}_'):
                                    matched_ocr_name = ocr_img_name
                                    found = True
                                    break
                    
                    # 如果文件名不匹配，尝试按大小匹配（旧格式兼容）
                    if not found:
                        for ocr_img_name, ocr_text in ocr_results.items():
                            if ocr_img_name in image_info_map:
                                _, actual_size = image_info_map[ocr_img_name]
                                if abs(actual_size - expected_size) < 100:
                                    matched_ocr_name = ocr_img_name
                                    found = True
                                    break
                    
                    if found and matched_ocr_name:
                        ocr_text = ocr_results[matched_ocr_name]
                        optimized_text = optimizer.optimize_image_ocr(ocr_text)
                        if optimized_text:
                            slide_ocr.append(f"**图片OCR [{matched_ocr_name}]**\n\n{optimized_text}")
                    elif image_name.lower().endswith('.emf'):
                        continue
                    else:
                        slide_ocr.append(f"**图片：{image_name}**\n\n（无文字内容）")
                
                if slide_ocr:
                    result_parts.append("\n\n---\n\n")
                    result_parts.append(f"**【幻灯片 {slide_num} - 图片识别内容】**\n\n")
                    result_parts.append("\n\n".join(slide_ocr))
                    result_parts.append("\n\n")
        
        elif slides[i].strip():
            result_parts.append(slides[i])
            i += 1
        else:
            i += 1

    # 合并结果
    final_content = ''.join(result_parts)
    
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(final_content)


def extract_pptx_to_markdown(pptx_path: str, output_dir: str = None, 
                              skip_ocr: bool = False, skip_semantic: bool = False):
    """
    Extract PPTX content to Markdown with image OCR.
    Output is placed in a folder with the same name as the PPTX file.
    """
    pptx_path = Path(pptx_path).resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    # 创建与PPT同名的输出文件夹
    if output_dir is None:
        output_dir = pptx_path.parent / pptx_path.stem
    else:
        output_dir = Path(output_dir) / pptx_path.stem
    
    # 确保输出目录存在
    output_dir.mkdir(parents=True, exist_ok=True)

    output_md = output_dir / f"{pptx_path.stem}_含图片内容.md"

    print(f"Processing: {pptx_path}")
    print(f"Output: {output_md}")

    install_dependencies()

    print("Extracting text content...")
    if not extract_markitdown(pptx_path, output_md):
        raise RuntimeError("Failed to extract text content")

    # 提取图片到输出目录（改进版：按幻灯片索引命名）
    print("Extracting images...")
    image_map, slide_to_images = extract_images_to_output(str(pptx_path), str(output_dir))
    print(f"Extracted {len(image_map)} images to 'images/' folder")
    if slide_to_images:
        print(f"Images mapped to {len(slide_to_images)} slides")

    # 修复Markdown中的图片链接（改进版：按幻灯片索引匹配）
    print("Fixing image links...")
    fix_image_links(str(output_md), image_map, slide_to_images)

    if not skip_ocr:
        # 获取幻灯片与图片的对应关系
        print("Mapping images to slides...")
        slide_images = get_slide_images(str(pptx_path))
        print(f"Found images in {len(slide_images)} slides")

        # 构建图片信息映射（用于OCR）
        image_info_map = {}
        for img_name, img_path in image_map.items():
            if os.path.exists(img_path):
                image_info_map[img_name] = (img_path, os.path.getsize(img_path))

        if image_info_map:
            print("Performing OCR on images...")
            ocr_results = {}
            for img_name, (img_path, _) in image_info_map.items():
                print(f"OCR: {img_name}")
                ocr_text = ocr_image(img_path)
                if ocr_text:
                    ocr_results[img_name] = ocr_text

            # 创建语义优化器
            optimizer = SemanticOptimizer()

            # 将OCR结果插入到对应的幻灯片位置
            print("Inserting OCR results into slides...")
            insert_ocr_results(str(output_md), ocr_results, slide_images, image_info_map, optimizer)
    else:
        print("OCR disabled")

    print(f"\nDone! Output: {output_md}")
    return output_md


def main():
    args = parse_args()

    try:
        extract_pptx_to_markdown(
            args.pptx_file,
            args.output_dir,
            args.no_ocr,
            args.no_semantic
        )
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()