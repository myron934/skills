"""Extract PPTX content to Markdown with image OCR.

Usage:
    python extract_pptx.py "path\to\file.pptx" [output_dir] [options]

Options:
    --no-ocr            Skip OCR for images
    -h, --help          Show this help message
"""

import os
import sys
import re
import zipfile
import shutil
import subprocess
import argparse
from pathlib import Path
from typing import Optional, Dict, List, Tuple

# 全局OCR Reader缓存
_ocr_reader = None


def parse_args():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(
        description="Extract PPTX content to Markdown with image OCR"
    )
    parser.add_argument("pptx_file", help="Path to the PPTX file")
    parser.add_argument("output_dir", nargs="?", default=None,
                        help="Output directory (default: same as input)")
    parser.add_argument("--no-ocr", action="store_true",
                        help="Skip OCR for images")
    return parser.parse_args()


def clean_ocr_text(text: str) -> str:
    """OCR文本清理：去除噪声、修正格式"""
    if not text:
        return text

    rules = [
        (r'\|\s*\|', ' '),           # 清理多余竖线
        (r' {2,}', ' '),             # 清理重复空格
        (r'([^\w\s\n\-=→]){2,}', ' '),  # 清理连续标点
    ]

    for pattern, replacement in rules:
        text = re.sub(pattern, replacement, text)

    return text.strip()


def is_noise_text(text: str) -> bool:
    """判断是否为无意义的噪声文本"""
    noise_patterns = [
        r'^\s*(TE|wh|we|CAVA|RR)\s*$',
        r'^\s*[DDEEAA\s.]+$',
        r'^\s*[|\\/\-=_~`^\*\+\>\<\?\[\]\(\)]+\s*$',
    ]
    return any(re.match(p, text, re.IGNORECASE) for p in noise_patterns)


def install_dependencies():
    """Ensure required packages are installed."""
    packages = ["markitdown[pptx]", "Pillow", "easyocr", "python-pptx"]
    print("Checking dependencies...")
    for pkg in packages:
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", pkg, "-q"],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            print(f"Warning: Failed to install {pkg}", file=sys.stderr)


def extract_markitdown(pptx_path: str, output_md: str) -> bool:
    """Extract PPTX text content using markitdown."""
    result = subprocess.run(
        [sys.executable, "-m", "markitdown", pptx_path],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        with open(output_md, 'w', encoding='utf-8') as f:
            f.write(result.stdout)
        return True
    print(f"markitdown error: {result.stderr}", file=sys.stderr)
    return False


def extract_images(pptx_path: str, output_dir: str) -> Tuple[Dict[str, dict], Dict[int, List[str]]]:
    """
    提取PPTX中的图片到输出目录
    返回: (图片信息映射, 幻灯片→图片列表映射)
    """
    images_dir = os.path.join(output_dir, "images")
    
    if os.path.exists(images_dir):
        shutil.rmtree(images_dir, ignore_errors=True)
    
    os.makedirs(images_dir, exist_ok=True)

    image_map = {}
    slide_to_images = {}

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_images = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                if not hasattr(shape, 'image'):
                    continue
                    
                image = shape.image
                original_name = image.filename
                if not original_name:
                    continue
                
                ext = os.path.splitext(original_name)[1].lower()
                if ext == '.emf':
                    continue
                
                image_data = image.blob
                output_name = f"slide_{slide_idx}_{shape_idx}{ext}"
                output_path = os.path.join(images_dir, output_name)
                
                with open(output_path, 'wb') as f:
                    f.write(image_data)
                
                image_map[output_name] = {
                    'path': output_path,
                    'size': len(image_data),
                    'slide_idx': slide_idx
                }
                slide_images.append(output_name)
            
            if slide_images:
                slide_to_images[slide_idx] = slide_images
    
    except ImportError:
        print("Warning: python-pptx not found, falling back to zip extraction", file=sys.stderr)
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for name in zf.namelist():
                if not name.startswith("ppt/media/"):
                    continue
                original_name = os.path.basename(name)
                if not original_name:
                    continue
                ext = os.path.splitext(original_name)[1].lower()
                if ext == '.emf':
                    continue
                output_path = os.path.join(images_dir, original_name)
                with open(output_path, 'wb') as f:
                    f.write(zf.read(name))
                image_map[original_name] = {
                    'path': output_path,
                    'size': os.path.getsize(output_path),
                    'slide_idx': None
                }
    
    return image_map, slide_to_images


def fix_image_links(md_path: str, slide_to_images: Dict[int, List[str]]):
    """修复Markdown中的图片链接"""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    lines = content.split('\n')
    result_lines = []
    current_slide = None
    slide_image_idx = {}

    for line in lines:
        slide_match = re.search(r'<!-- Slide number: (\d+) -->', line)
        if slide_match:
            current_slide = int(slide_match.group(1))
            slide_image_idx[current_slide] = 0
            result_lines.append(line)
        elif current_slide and current_slide in slide_to_images:
            img_match = re.search(r'!\[([^\]]*)\]\(([^)]+)\)', line)
            if img_match:
                alt_text = img_match.group(1)
                images = slide_to_images[current_slide]
                idx = slide_image_idx.get(current_slide, 0)
                
                if idx < len(images):
                    image_name = images[idx]
                    slide_image_idx[current_slide] = idx + 1
                    line = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', 
                                  f'![{alt_text}](images/{image_name})', line)
            result_lines.append(line)
        else:
            result_lines.append(line)

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(result_lines))


def get_ocr_reader():
    """获取或初始化OCR Reader（缓存）"""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(['ch_sim', 'en'], verbose=False)
    return _ocr_reader


def ocr_image(img_path: str) -> Optional[str]:
    """Perform OCR on a single image."""
    try:
        import numpy as np
        import cv2

        ext = os.path.splitext(img_path)[1].lower()
        if ext == '.emf':
            return None

        # 使用numpy读取图片（解决OpenCV不支持中文路径的问题）
        with open(img_path, 'rb') as f:
            img_array = np.frombuffer(f.read(), dtype=np.uint8)
        img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
        
        if img is None:
            return None

        reader = get_ocr_reader()
        result = reader.readtext(img)
        
        text = ' '.join([item[1] for item in result])
        if text.strip():
            return clean_ocr_text(text.strip())
        return None

    except ImportError:
        print("Warning: easyocr/numpy/cv2 not found", file=sys.stderr)
        return None
    except Exception as e:
        print(f"OCR error for {img_path}: {e}", file=sys.stderr)
        return None


def insert_ocr_results(md_path: str, ocr_results: Dict[str, str], 
                       slide_to_images: Dict[int, List[str]]):
    """将OCR结果插入到对应的幻灯片位置"""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = re.split(r'(<!-- Slide number: (\d+) -->)', content)
    result_parts = []
    i = 0
    
    while i < len(slides):
        if i + 2 < len(slides) and slides[i].startswith('<!-- Slide number:'):
            slide_num = int(slides[i+1])
            result_parts.append(slides[i])
            result_parts.append(slides[i+1])
            i += 2
            if i < len(slides):
                result_parts.append(slides[i])
                i += 1
            
            if slide_num in slide_to_images:
                slide_ocr = []
                for image_name in slide_to_images[slide_num]:
                    if image_name in ocr_results:
                        text = ocr_results[image_name]
                        if not is_noise_text(text):
                            slide_ocr.append(f"**图片OCR [{image_name}]**\n\n{text}")
                    else:
                        slide_ocr.append(f"**图片：{image_name}**\n\n（无文字内容）")
                
                if slide_ocr:
                    result_parts.append("\n\n---\n\n")
                    result_parts.append(f"**【幻灯片 {slide_num} - 图片识别内容】**\n\n")
                    result_parts.append("\n\n".join(slide_ocr))
                    result_parts.append("\n\n")
        
        elif slides[i].strip():
            result_parts.append(slides[i])
            i += 1
        else:
            i += 1

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(''.join(result_parts))


def extract_pptx_to_markdown(pptx_path: str, output_dir: str = None, skip_ocr: bool = False):
    """Extract PPTX content to Markdown with image OCR."""
    pptx_path = Path(pptx_path).resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if output_dir is None:
        output_dir = pptx_path.parent / pptx_path.stem
    else:
        output_dir = Path(output_dir) / pptx_path.stem
    
    output_dir.mkdir(parents=True, exist_ok=True)
    output_md = output_dir / f"{pptx_path.stem}_含图片内容.md"

    print(f"Processing: {pptx_path}")
    print(f"Output: {output_md}")

    install_dependencies()

    print("Extracting text content...")
    if not extract_markitdown(str(pptx_path), str(output_md)):
        raise RuntimeError("Failed to extract text content")

    print("Extracting images...")
    image_map, slide_to_images = extract_images(str(pptx_path), str(output_dir))
    print(f"Extracted {len(image_map)} images to 'images/' folder")

    print("Fixing image links...")
    fix_image_links(str(output_md), slide_to_images)

    if not skip_ocr and image_map:
        print("Performing OCR on images...")
        ocr_results = {}
        for output_name, info in image_map.items():
            print(f"OCR: {output_name}")
            text = ocr_image(info['path'])
            if text:
                ocr_results[output_name] = text

        print("Inserting OCR results...")
        insert_ocr_results(str(output_md), ocr_results, slide_to_images)

    print(f"\nDone! Output: {output_md}")
    return output_md


def main():
    args = parse_args()
    try:
        extract_pptx_to_markdown(args.pptx_file, args.output_dir, args.no_ocr)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
